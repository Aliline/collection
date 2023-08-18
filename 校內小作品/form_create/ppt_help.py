# -*- coding: UTF-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.dml.fill import FillFormat
prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

left = top = width = height = Inches(1)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame


p = tf.add_paragraph()
p.text = "日月光"
# 大小
p.font.size = Pt(18)
# 字形
p.font.name = "Microsoft YaHei"

p = tf.add_paragraph()
p.text = "日月光"
p.font.size = Pt(18)
# 粗體
p.font.bold = True
p.font.name = "Microsoft YaHei"

p = tf.add_paragraph()
p.text = "日月光"
p.font.size = Pt(18)
# 斜體
p.font.italic = True
p.font.name = "Microsoft YaHei"


p = tf.add_paragraph()
p.text = "日月光"
p.font.size = Pt(18)
# 底線
p.font.underline = True
p.font.name = "Microsoft YaHei"

p = tf.add_paragraph()
p.text = "日月光"
p.font.size = Pt(18)
# 刪除線
p.font._element.attrib['strike'] = 'sngStrike'
p.font.name = "Microsoft YaHei"


prs.save('test.pptx')
